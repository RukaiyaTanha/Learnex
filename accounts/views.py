from django.shortcuts import render
from django.http import JsonResponse,FileResponse
from django.contrib.auth import get_user_model, authenticate, login, logout
from django.views.decorators.csrf import csrf_exempt
from django.contrib.auth.decorators import login_required
from .models import Course, UserCourse,Topic, Question,Marks,Syllabus, QuizAttendance,TeacherCourseSelection,StudentInfo,Section, TeacherSyllabusUpload, AIQuizAttempt
from django.contrib.auth.tokens import default_token_generator
from django.utils.http import urlsafe_base64_encode, urlsafe_base64_decode
from django.utils.encoding import force_bytes
from django.core.mail import send_mail
from django.template.loader import render_to_string
from django.shortcuts import render, redirect, get_object_or_404
from django.core.mail import EmailMultiAlternatives
from .models import CustomUser
from django.urls import reverse
import json as pyjson
from difflib import SequenceMatcher
from django.contrib import messages
from django.views.decorators.http import require_POST
import mimetypes
from django.db.models import F, Avg
import csv, io, json
import os
import pdfplumber
from PIL import Image
import pytesseract
import docx
from pptx import Presentation
import openpyxl
from io import BytesIO
from django.conf import settings
from django.contrib import messages
import google.generativeai as genai
import markdown2
from .ai_feedback import generate_ai_feedback
from django.utils.safestring import mark_safe
from datetime import datetime
from django.db import transaction
from collections import defaultdict
from datetime import timedelta
from django.utils.timezone import now

User = get_user_model()

# ---------- Registration API ----------
@csrf_exempt
def register(request):
    if request.method == "POST":
        data = json.loads(request.body)
        email = data.get('email')
        password = data.get('password')
        role = data.get('role')
        username = data.get('name')
        first_name = data.get('first_name')
        last_name = data.get('last_name')

        if not email or not password or not role or not username:
            return JsonResponse({"success": False, "message": "Missing fields"})

        if User.objects.filter(email=email).exists():
            return JsonResponse({"success": False, "message": "Email already registered"})

        user = User.objects.create_user(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            password=password,
            role=role
        )
        return JsonResponse({"success": True, "message": "User registered successfully"})

    return JsonResponse({"success": False, "message": "Invalid request"})


# ---------- Registration HTML page ----------
def registration_page(request):
    return render(request, "accounts/registration.html")


# ---------- Login API ----------
@csrf_exempt
def login_api(request):
    if request.method == "POST":
        data = json.loads(request.body)
        email = data.get("email")
        password = data.get("password")
        role = data.get("role")  

        if not email or not password or not role:
            return JsonResponse({"success": False, "message": "Missing email, password, or role"})

        try:
            user = User.objects.get(email=email)

            if not user.check_password(password):
                return JsonResponse({"success": False, "message": "Invalid email or password"})

            if getattr(user, "role", "").lower() != role.lower():
                return JsonResponse({"success": False, "message": f"This account is not a {role}"})

            login(request, user)

            if role.lower() == "student":
                redirect_url = reverse("student_dashboard")
            elif role.lower() == "teacher":
                redirect_url = reverse("teacher_dashboard")
            else:
                redirect_url = reverse("student_dashboard")

            return JsonResponse({
                "success": True,
                "message": "Login successful",
                "username": user.username,
                "role": getattr(user, "role", "N/A"),
                "redirect_url": redirect_url,
            })

        except User.DoesNotExist:
            return JsonResponse({"success": False, "message": "Invalid email or password"})

    return JsonResponse({"success": False, "message": "Invalid request"})

# ---------- Login HTML page ----------
def login_page(request):
    return render(request, "accounts/login.html")

# Student Dashboard page (protected)
@login_required(login_url='/accounts/login-page/')
def student_dashboard(request):
    user = request.user
    context = {
        "username": user.username,
        "role": getattr(user, "role", "N/A"),
    }
    return render(request, "accounts/dashboard.html", context)

# Logout API
@csrf_exempt
def logout_api(request):
    if request.method == "POST":
        logout(request)
        return JsonResponse({"success": True, "message": "Logged out"})
    return JsonResponse({"success": False, "message": "Invalid request"})

#save_courses_api
@csrf_exempt
@login_required(login_url='/accounts/login-page/')
def save_courses_api(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            selected_codes = data.get("courses", [])
            if not isinstance(selected_codes, list):
                return JsonResponse({"success": False, "message": "Invalid courses data"})

            user = request.user

            all_courses = Course.objects.filter(code__in=selected_codes)

            added = []
            for course in all_courses:
                uc, created = UserCourse.objects.get_or_create(user=user, course=course)
                if created:
                    added.append(course.code)

            removed = []
            user_courses_qs = UserCourse.objects.filter(user=user)
            for uc in user_courses_qs:
                if uc.course.code not in selected_codes:
                    removed.append(uc.course.code)
                    uc.delete()

            current = list(UserCourse.objects.filter(user=user).values_list('course__code', flat=True))

            return JsonResponse({
                "success": True,
                "message": f"Added {len(added)} course(s), removed {len(removed)} course(s).",
                "courses": current
            })
        except Exception as e:
            return JsonResponse({"success": False, "message": str(e)})

    return JsonResponse({"success": False, "message": "Invalid request"})

@login_required(login_url='/accounts/login-page/')
def course_selection_page(request):
    user = request.user
    courses = Course.objects.all()
    user_courses_qs = UserCourse.objects.filter(user=user).values_list('course__code', flat=True)
    user_courses = list(user_courses_qs)

    context = {
        "courses": courses,
        "user_courses": user_courses,                    
        "user_courses_json": pyjson.dumps(user_courses),  
        "username": user.username,
    }
    return render(request, "accounts/course_selection.html", context)

@login_required(login_url='/accounts/login-page/')
def selected_courses_page(request):
    user_courses = UserCourse.objects.filter(user=request.user).select_related('course')
    return render(request, 'accounts/selected_courses.html', {
        'username': request.user.username,
        'user_courses': user_courses,
    })

# ---------- Quiz Selection ----------
@login_required(login_url='/accounts/login-page/')
def quiz_selection_page(request):
    user = request.user
    enrolled_courses = Course.objects.filter(usercourse__user=user)

    course_topics = {str(course.id): [{"id": t.id, "name": t.name} for t in course.topics.all()] for course in enrolled_courses}

    context = {
        "username": user.username,
        "courses": enrolled_courses,
        "course_topics": course_topics,
    }
    return render(request, "accounts/quiz_selection.html", context)

# ---------- Start Quiz ----------
@login_required(login_url='/accounts/login-page/')
def start_quiz(request, course_id, topic_id, quiz_type):
    try:
        course = Course.objects.get(id=course_id)
        topic = Topic.objects.get(id=topic_id, course=course)
    except (Course.DoesNotExist, Topic.DoesNotExist):
        return render(request, "accounts/quiz_invalid.html", {"message": "Invalid course or topic."})

    questions = Question.objects.filter(course=course, topic=topic, q_type=quiz_type)

    context = {
        "username": request.user.username,
        "course": course,
        "topic": topic,
        "type": quiz_type,
        "questions": questions,
    }
    return render(request, "accounts/quiz_page.html", context)

# ---------- Submit Quiz ----------
@login_required(login_url='/accounts/login-page/')
def submit_quiz(request):
    if request.method == "POST":
        score = 0
        total = 0
        for key, value in request.POST.items():
            if key.startswith("q_"):
                qid = int(key.split("_")[1])
                try:
                    question = Question.objects.get(id=qid)
                    total += 1
                    if question.answer == value:
                        score += 1
                except Question.DoesNotExist:
                    continue
        return render(request, "accounts/quiz_result.html", {"score": score, "total": total})

# Forgot Password API
@csrf_exempt
def forgot_password_api(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        email = data.get('email')

        try:
            user = CustomUser.objects.get(email=email)
        except CustomUser.DoesNotExist:
            return JsonResponse({'success': False, 'message': 'No account found with that email.'})

        uidb64 = urlsafe_base64_encode(force_bytes(user.pk))
        token = default_token_generator.make_token(user)
        reset_link = request.build_absolute_uri(
            reverse('reset_password_page', kwargs={'uidb64': uidb64, 'token': token})
        )

        html_content = render_to_string('accounts/reset_password_email.html', {
            'username': user.username,
            'reset_link': reset_link
        })
        text_content = f"Hi {user.username},\n\nPlease click the link to reset your password:\n{reset_link}"

        email_message = EmailMultiAlternatives(
            subject='Smart Class Password Reset',
            body=text_content,
            from_email='smartclass@example.com',
            to=[user.email],
        )
        email_message.attach_alternative(html_content, "text/html")
        email_message.send()

        return JsonResponse({'success': True, 'message': 'Password reset email sent!'})

# Forgot Password Page
def forgot_password_page(request):
    return render(request, 'accounts/forgot_password.html')

# Reset Password Page
def reset_password_page(request, uidb64, token):
    try:
        uid = urlsafe_base64_decode(uidb64).decode()
        user = CustomUser.objects.get(pk=uid)
    except (TypeError, ValueError, OverflowError, CustomUser.DoesNotExist):
        user = None

    if user and default_token_generator.check_token(user, token):
        return render(request, "accounts/reset_password.html", {"uidb64": uidb64, "token": token})
    else:
        return render(request, "accounts/reset_invalid.html")

# Reset Password API
@csrf_exempt
def reset_password_api(request):
    if request.method == "POST":
        data = json.loads(request.body)
        uidb64 = data.get("uidb64")
        token = data.get("token")
        new_password = data.get("password")

        try:
            uid = urlsafe_base64_decode(uidb64).decode()
            user = CustomUser.objects.get(pk=uid)
        except (TypeError, ValueError, OverflowError, CustomUser.DoesNotExist):
            return JsonResponse({"success": False, "message": "Invalid link"})

        if user and default_token_generator.check_token(user, token):
            user.set_password(new_password)
            user.save()
            return JsonResponse({"success": True, "message": "Password reset successful"})
        else:
            return JsonResponse({"success": False, "message": "Invalid or expired link"})

    return JsonResponse({"success": False, "message": "Invalid request"})

@login_required(login_url='/accounts/login-page/')
def profile_page(request):
    user = request.user
    context = {
        "username": user.username,
        "first_name": user.first_name,
        "last_name": user.last_name,
        "role": getattr(user, "role", "N/A"),
        "email": user.email,
        "date_joined": user.date_joined,
    }
    return render(request, "accounts/profile.html", context)

# ---------- Edit Profile Page ----------
@login_required(login_url='/accounts/login-page/')
def profile_edit(request):
    user = request.user
    context = {
        "username": user.username,
        "first_name": user.first_name,
        "last_name": user.last_name,
        "email": user.email,
        "role": getattr(user, "role", "N/A"),
    }
    return render(request, "accounts/profile_edit.html", context)

# ---------- Update Profile API ----------
@csrf_exempt
@login_required(login_url='/accounts/login-page/')
def update_profile_api(request):
    if request.method == "POST":
        try:
            data = request.POST
            profile_picture = request.FILES.get('profile_picture')
            user = request.user

            username = data.get("username", "").strip()
            first_name = data.get("first_name", "").strip()
            last_name = data.get("last_name", "").strip()
            email = data.get("email", "").strip()
            profile_picture = request.FILES.get("profile_picture")

            if not username or not email or not first_name or not last_name:
                return JsonResponse({"success": False, "message": "All fields are required."})

            if User.objects.filter(email=email).exclude(id=user.id).exists():
                return JsonResponse({"success": False, "message": "This email is already in use."})

            user.username = username
            user.first_name = first_name
            user.last_name = last_name
            user.email = email
            if profile_picture:
                user.profile_picture = profile_picture
            user.save()

            return JsonResponse({"success": True, "message": "Profile updated successfully!"})
        except Exception as e:
            return JsonResponse({"success": False, "message": str(e)})

    return JsonResponse({"success": False, "message": "Invalid request"})

# ---------- Quiz page ----------
@login_required(login_url='/accounts/login-page/')
def quiz_page(request, course_id, topic_id, type):
    try:
        course = Course.objects.get(id=course_id)
        topic = Topic.objects.get(id=topic_id, course=course)
        topic_name = topic.name
    except (Course.DoesNotExist, Topic.DoesNotExist):
        return render(request, "accounts/quiz_invalid.html", {"message": "Invalid course or topic selected."})

    all_questions = {
    "6": {  # Course: Differential
        "94": {  # Topic: Limits
            "mcq": {
                "time_duration": "15 Minutes",
                "marks": 20,
                "questions": [
                    {"question": "lim x→2 (x^2 - 4)/(x - 2)", "options": ["2", "4", "6", "8"], "answer": "4"},
                    {"question": "lim x→0 sin(x)/x", "options": ["0", "1", "Infinity", "-Infinity"], "answer": "1"},
                    {"question": "lim x→1 (x^3 - 1)/(x - 1)", "options": ["1", "2", "3", "0"], "answer": "3"},
                    {"question": "lim x→0 |x|/x", "options": ["0", "1", "Does not exist", "-1"], "answer": "Does not exist"},
                    {"question": "If lim x→a f(x)=L, it may or may not equal f(a)", "options": ["True", "False"], "answer": "True"},
                    {"question": "lim x→∞ 1/x", "options": ["0", "1", "Infinity", "-Infinity"], "answer": "0"},
                    {"question": "lim x→0 (1+x)^(1/x)", "options": ["1", "e", "0", "Infinity"], "answer": "e"},
                    {"question": "Rule used to solve 0/0 form", "options": ["L’Hôpital’s Rule", "Chain Rule", "Product Rule", "Quotient Rule"], "answer": "L’Hôpital’s Rule"},
                    {"question": "lim x→0 (1-cos(x))/x^2", "options": ["0", "1/2", "1", "2"], "answer": "1/2"},
                    {"question": "lim x→0 tan(x)/x", "options": ["0", "1", "Infinity", "-Infinity"], "answer": "1"},
                    {"question": "lim x→3 (x^2+2x-9)/(x-3)", "options": ["6", "7", "8", "9"], "answer": "8"},
                    {"question": "lim x→0 (e^x - 1)/x", "options": ["0", "1", "e", "-1"], "answer": "1"},
                    {"question": "If left-hand ≠ right-hand limits, limit does not exist", "options": ["True", "False"], "answer": "True"},
                    {"question": "lim x→0 x*ln(x)", "options": ["0", "1", "-Infinity", "Does not exist"], "answer": "0"},
                    {"question": "lim x→∞ (1 + 2/x)^x", "options": ["2", "e", "e^2", "Infinity"], "answer": "e^2"},
                    {"question": "For f(x)=|x|, lim x→0^- f(x)", "options": ["0", "1", "-1", "Does not exist"], "answer": "0"},
                    {"question": "Removable discontinuity is called?", "options": ["Hole", "Jump", "Infinite", "Oscillation"], "answer": "Hole"},
                    {"question": "lim x→0 sin(3x)/x", "options": ["1", "3", "0", "Does not exist"], "answer": "3"},
                    {"question": "lim x→∞ (2x^2+5x)/(x^2+3x)", "options": ["1", "2", "3", "5"], "answer": "2"},
                    {"question": "Formal definition of limit involves |x-a|<δ ⇒ |f(x)-L|<ε", "options": ["True", "False"], "answer": "True"}
                ]
            },
            "scenario-mcq": {
                "time_duration": "20 Minutes",
                "marks": 20,
                "questions": [
                    {"question": "Function jumps from 2 to 5 at x=1", "options": ["Limit exists", "Limit does not exist"], "answer": "Limit does not exist"},
                    {"question": "LHL=RHL=2", "options": ["Limit exists", "Limit does not exist"], "answer": "Limit exists"},
                    {"question": "Graph has hole at (1,2)", "options": ["Limit = 0", "Limit = 2"], "answer": "Limit = 2"},
                    {"question": "lim Δt→0 Δs/Δt", "options": ["Displacement", "Velocity", "Acceleration", "Speed"], "answer": "Velocity"},
                    {"question": "P(t)=100t/(t+1), t→∞", "options": ["100", "0", "Infinity", "1"], "answer": "100"},
                    {"question": "E(x)=(x^2-4)/(x-2)", "options": ["Exact error", "Approximated error"], "answer": "Approximated error"},
                    {"question": "T(t)=5+2*sin(t), t→0", "options": ["5", "0", "2", "Infinity"], "answer": "5"},
                    {"question": "f(x)=1/x, x→0+", "options": ["0", "∞", "-∞", "Does not exist"], "answer": "∞"},
                    {"question": "sin(0.00001)/0.00001", "options": ["≈0", "≈1", "≈2", "Undefined"], "answer": "≈1"},
                    {"question": "Limit exists but f(a) undefined", "options": ["Function continuous", "Function discontinuous"], "answer": "Function discontinuous"},
                    {"question": "Infinite oscillation (sin(1/x))", "options": ["Limit exists", "Limit does not exist"], "answer": "Limit does not exist"},
                    {"question": "Growth rate→0", "options": ["Growth continues", "Growth stops"], "answer": "Growth stops"},
                    {"question": "Left and right limits differ", "options": ["Bilateral limit", "Jump discontinuity"], "answer": "Jump discontinuity"},
                    {"question": "Supply→∞, price→0", "options": ["Direct relation", "Inverse relation"], "answer": "Inverse relation"},
                    {"question": "f(x)=1/(x-2), x→2", "options": ["Limit exists", "Limit does not exist"], "answer": "Limit does not exist"},
                    {"question": "Same value both sides", "options": ["Bilateral limit", "Jump discontinuity"], "answer": "Bilateral limit"},
                    {"question": "T(x)=x^2, x→-3", "options": ["Limit = -9", "Limit = 9"], "answer": "Limit = 9"},
                    {"question": "f(x)=|x|/x, x=0", "options": ["Defined", "Undefined"], "answer": "Undefined"},
                    {"question": "(1+1/x)^x, x→∞", "options": ["1", "e"], "answer": "e"},
                    {"question": "Limit depends on path", "options": ["Path-dependent", "Path-independent"], "answer": "Path-dependent"}
                ]
            },
            "code": {
                "time_duration": "1 Hour",
                "marks": 35,
                "questions": [
                    {"question": "Python program to approximate limit of sin(x)/x as x→0",  "options": [],"answer": "import math\nx = 0.000001\nprint(math.sin(x)/x)"},
                    {"question": "Python program to calculate (1 + 1/x)^x as x→∞", "options": [], "answer": "x = 10000\nprint((1 + 1/x) ** x)"},
                    {"question": "Trace output:\nimport math\nx = 0.000001\nprint(math.sin(x)/x)",  "options": [], "answer": "0.9999999999998334"},
                    {"question": "Trace output:\nx = 10000\nprint((1 + 1/x)**x)", "options": [], "answer": "2.7181459"},
                    {"question": "Trace output:\ndef f(x):\n    return (x**2 - 4)/(x - 2)\nprint(f(2.0001))",  "options": [],"answer": "4.0001"}
                ]
            },
            "theory": {
                "time_duration": "45 Minutes",
                "marks": 30,
                "questions": [
                    {"question": "What is a limit?", "options": [], "answer": "The value a function approaches as the input approaches a specific point."},
                    {"question": "What are left-hand and right-hand limits?",  "options": [],"answer": "Left-hand (x→a⁻) approaches from left; Right-hand (x→a⁺) from right. Limit exists if both equal."},
                    {"question": "What is an indeterminate form?", "options": [], "answer": "A form like 0/0, ∞/∞, 0×∞, where direct substitution fails."},
                    {"question": "Explain continuity using limits.",  "options": [],"answer": "A function is continuous at x=a if lim x→a f(x) = f(a)."},
                    {"question": "What is the geometric meaning of a limit?", "options": [], "answer": "The point the function’s graph approaches as x gets closer to a value."},
                    {"question": "Discuss types of discontinuities with examples.",  "options": [],"answer": "Removable: Limit exists but function value missing or different. Jump: Left and right limits unequal. Infinite: Function tends to ∞ near a point."},
                    {"question": "Explain L’Hôpital’s Rule with examples.", "options": [], "answer": "If a limit gives 0/0 or ∞/∞, lim x→a f(x)/g(x) = lim x→a f'(x)/g'(x). Examples: lim x→0 sin(x)/x = 1; lim x→∞ ln(x)/x = 0."}
                ]
            }
        }
    },
    "1": {  # Course: Web Technology
       "84": {  # Topic: HTML Basics
            "mcq": {
            "time_duration": "15 Minutes",
            "marks": 20,
            "questions": [
                {"question": "Which tag is used for paragraph?", "options": ["<p>", "<div>", "<span>", "<h1>"], "answer": "<p>"},
                {"question": "Which tag is used for the largest heading?", "options": ["<h1>", "<h6>", "<head>", "<header>"], "answer": "<h1>"},
                {"question": "Which tag is used to create a hyperlink?", "options": ["<a>", "<link>", "<url>", "<button>"], "answer": "<a>"},
                {"question": "Which attribute specifies the link URL in <a>?", "options": ["href", "src", "link", "url"], "answer": "href"},
                {"question": "Which tag is used for unordered list?", "options": ["<ul>", "<ol>", "<li>", "<list>"], "answer": "<ul>"},
                {"question": "Which tag is used for ordered list?", "options": ["<ol>", "<ul>", "<li>", "<list>"], "answer": "<ol>"},
                {"question": "Which tag is used for table row?", "options": ["<tr>", "<td>", "<table>", "<th>"], "answer": "<tr>"},
                {"question": "Which tag is used for table header?", "options": ["<th>", "<td>", "<tr>", "<thead>"], "answer": "<th>"},
                {"question": "Which tag is used to embed image?", "options": ["<img>", "<picture>", "<image>", "<src>"], "answer": "<img>"},
                {"question": "Which attribute is used for alternative text for images?", "options": ["alt", "title", "src", "href"], "answer": "alt"},
                {"question": "Which tag is used for bold text?", "options": ["<b>", "<strong>", "<bold>", "<em>"], "answer": "<b>"},
                {"question": "Which tag is used for italic text?", "options": ["<i>", "<em>", "<italic>", "<it>"], "answer": "<i>"},
                {"question": "Which tag defines a division or section?", "options": ["<div>", "<span>", "<section>", "<article>"], "answer": "<div>"},
                {"question": "Which tag is inline element?", "options": ["<span>", "<div>", "<p>", "<header>"], "answer": "<span>"},
                {"question": "Which tag is block-level element?", "options": ["<div>", "<span>", "<a>", "<img>"], "answer": "<div>"},
                {"question": "Which tag is used for adding page title?", "options": ["<title>", "<head>", "<meta>", "<h1>"], "answer": "<title>"},
                {"question": "Which tag is used to add a comment?", "options": ["<!-- -->", "//", "#", "/* */"], "answer": "<!-- -->"},
                {"question": "Which tag is used for line break?", "options": ["<br>", "<hr>", "<lb>", "<break>"], "answer": "<br>"},
                {"question": "Which tag is used to define a list item?", "options": ["<li>", "<ul>", "<ol>", "<list>"], "answer": "<li>"},
                {"question": "Which tag is used for embedding video?", "options": ["<video>", "<media>", "<embed>", "<movie>"], "answer": "<video>"}
            ]
        },
        "scenario-mcq": {
            "time_duration": "20 Minutes",
            "marks": 20,
            "questions": [
                {"question": "You want a navigation menu with links to Home, About, Contact. Which tag should you use?", "options": ["<nav>", "<div>", "<ul>", "<header>"], "answer": "<nav>"},
                {"question": "You want a table with header and 3 rows of data. Which tags are necessary?", "options": ["<table>, <tr>, <th>, <td>", "<table>, <row>, <th>, <td>", "<table>, <tr>, <td>", "<table>, <tr>, <header>, <cell>"], "answer": "<table>, <tr>, <th>, <td>"},
                {"question": "You want to display an image logo.png with alt text. Which tag and attribute?", "options": ["<img src='logo.png' alt='Logo'>", "<image src='logo.png' alt='Logo'>", "<picture src='logo.png'>", "<div src='logo.png'>"], "answer": "<img src='logo.png' alt='Logo'>"},
                {"question": "You want text 'Important' to be emphasized. Which tag is correct?", "options": ["<strong>", "<b>", "<em>", "<i>"], "answer": "<strong>"},
                {"question": "You need a horizontal line to separate sections. Which tag?", "options": ["<hr>", "<br>", "<line>", "<divider>"], "answer": "<hr>"},
                {"question": "You want a form input for user email. Which tag?", "options": ["<input type='email'>", "<input type='text'>", "<form type='email'>", "<email>"], "answer": "<input type='email'>"},
                {"question": "You want to embed a YouTube video. Which tag?", "options": ["<iframe>", "<embed>", "<video>", "<object>"], "answer": "<iframe>"},
                {"question": "You want text aligned center. Which attribute/tag?", "options": ["<p style='text-align:center'>", "<center>", "<div align='center'>", "All of the above"], "answer": "All of the above"},
                {"question": "You want a dropdown selection in a form. Which tag?", "options": ["<select>", "<input>", "<option>", "<dropdown>"], "answer": "<select>"},
                {"question": "You want multiple checkboxes in a form. Which tag?", "options": ["<input type='checkbox'>", "<input type='radio'>", "<checkbox>", "<check>"], "answer": "<input type='checkbox'>"},
                {"question": "You want a numbered list. Which tag?", "options": ["<ol>", "<ul>", "<li>", "<list>"], "answer": "<ol>"},
                {"question": "You want a tooltip to appear on hover. Which attribute?", "options": ["title", "alt", "hover", "tooltip"], "answer": "title"},
                {"question": "You want a footer section with copyright info. Which tag?", "options": ["<footer>", "<bottom>", "<section>", "<div>"], "answer": "<footer>"},
                {"question": "You want a sidebar menu. Which semantic tag?", "options": ["<aside>", "<section>", "<nav>", "<div>"], "answer": "<aside>"},
                {"question": "You want to highlight text. Which tag?", "options": ["<mark>", "<highlight>", "<em>", "<strong>"], "answer": "<mark>"},
                {"question": "You want to group multiple form elements. Which tag?", "options": ["<fieldset>", "<form>", "<div>", "<section>"], "answer": "<fieldset>"},
                {"question": "You want to create a clickable image link. Which combination?", "options": ["<a><img></a>", "<img><a></img>", "<link><img></link>", "<div><img></div>"], "answer": "<a><img></a>"},
                {"question": "You want a semantic article block for blog content. Which tag?", "options": ["<article>", "<section>", "<div>", "<main>"], "answer": "<article>"},
                {"question": "You want to create a form with submit button. Which tag?", "options": ["<button>", "<input type='submit'>", "Both", "<form>"], "answer": "Both"},
                {"question": "You want to include external CSS file. Which tag?", "options": ["<link>", "<style>", "<css>", "<import>"], "answer": "<link>"}
            ]
        },
        "code": {
            "time_duration": "1 Hour",
            "marks": 35,
            "questions": [
                {"question": "Write HTML code to create a hyperlink to https://example.com with text 'Visit Example'", "options": [], "answer": "<a href='https://example.com'>Visit Example</a>"},
                {"question": "Write HTML code to create an ordered list with 3 items: Apple, Banana, Cherry", "options": [], "answer": "<ol><li>Apple</li><li>Banana</li><li>Cherry</li></ol>"},
                {"question": "Trace the output: <ul><li>One</li><li>Two</li></ul> (How many bullets will appear?)", "options": [], "answer": "2 bullets"},
                {"question": "Trace the output: <p style='text-align:center'>Hello</p>", "options": [], "answer": "Text 'Hello' centered"},
                {"question": "Trace the output: <img src='logo.png' alt='Logo'>", "options": [], "answer": "Displays image 'logo.png' with alt text 'Logo'"}
            ]
        },
        "theory": {
            "time_duration": "45 Minutes",
            "marks": 30,
            "questions": [
                {"question": "Define HTML and its purpose in web development.", "options": [], "answer": "HTML (HyperText Markup Language) is used to structure content on the web using elements and tags."},
                {"question": "Explain the difference between block-level and inline elements.", "options": [], "answer": "Block-level elements start on a new line and take full width (<div>, <p>), while inline elements do not start on a new line and take only needed width (<span>, <a>)."},
                {"question": "Describe the structure of an HTML document with its main sections.", "options": [], "answer": "An HTML document has <!DOCTYPE html>, <html>, <head> (meta info, title), <body> (visible content)."},
                {"question": "Explain semantic HTML and provide examples.", "options": [], "answer": "Semantic HTML elements describe meaning: <header>, <footer>, <article>, <section>."},
                {"question": "Explain the use of attributes in HTML tags with examples.", "options": [], "answer": "Attributes provide additional info about elements, e.g., <img src='logo.png' alt='Logo'> where src and alt are attributes."},
                {"question": "Discuss the difference between <div> and <section>.", "options": [], "answer": "<div> is generic container, <section> is semantic section of related content."},
                {"question": "Explain how HTML forms work and their main elements.", "options": [], "answer": "Forms collect user input. Main elements: <form>, <input>, <textarea>, <button>, <select>, <label>."}
            ]
        }
    }
}
}

    quiz_data = all_questions.get(str(course_id), {}).get(str(topic_id), {}).get(type, {})
    questions = quiz_data.get("questions", [])
    time_duration = quiz_data.get("time_duration", "N/A")
    marks = quiz_data.get("marks", "N/A")

    if request.method == "POST":
        user_answers = {}
        correct_count = 0
        results = []

        for i, q in enumerate(questions, start=1):
            selected = request.POST.get(f'answer_{i}')
            correct = q["answer"]
            is_correct = (selected == correct)

            user_answers[i] = selected
            if is_correct:
                correct_count += 1

            results.append({
                "question": q["question"],
                "options": q.get("options", []),
                "selected": selected,
                "correct": correct,
                "is_correct": is_correct,
            })

        total_score = correct_count
        performance_rate = round((correct_count / len(questions)) * 100, 2) if questions else 0

        ai_feedback = ""
        if type in ["mcq","scenario-mcq","code", "theory"]:
            try:
                ai_feedback = generate_ai_feedback(performance_rate, results, topic.name)
            except Exception as e:
                ai_feedback = f"(Feedback unavailable — {str(e)})"

        week_number = datetime.now().isocalendar()[1] % 4  
        if week_number == 0:
            week_number = 4
            
        QuizAttendance.objects.update_or_create(
            student=request.user,
            course=course,
            week=week_number,
            month=datetime.now().month, 
            defaults={
                "attended": True,
                "date": datetime.now()
                }
        )    
        
        context = {
            "course": course,
            "topic": topic,
            "type": type,
            "total_score": total_score,
            "performance_rate": performance_rate,
            "results": results,
            "ai_feedback": ai_feedback,
        }

        return render(request, "accounts/quiz_result.html", context)

    context = {
        "username": request.user.username,
        "course": course,
        "topic": topic,
        "type": type,
        "questions": questions,
        "time_duration": time_duration,
        "marks": marks,
    }
    return render(request, "accounts/quiz_page.html", context)

@login_required(login_url='/accounts/login-page/')
def upload_marks_page(request):
    user = request.user
    courses = Course.objects.filter(usercourse__user=user)  
    return render(request, "accounts/upload_marks.html", {"username": user.username, "courses": courses})


@login_required(login_url='/accounts/login-page/')
@csrf_exempt
def upload_marks_api(request):
    if request.method == "POST":
        data = json.loads(request.body)
        student = request.user 
        course = Course.objects.get(id=data['course_id'])

        marks, created = Marks.objects.update_or_create(
            student=student,
            course=course,
            defaults={
                "quiz1": data.get("quiz1", 0),
                "quiz2": data.get("quiz2", 0),
                "quiz3": data.get("quiz3", 0),
                "attendance": data.get("attendance", 0),
                "assignment": data.get("assignment", 0),
                "presentation": data.get("presentation", 0),
                "termexam": data.get("termexam", 0),
            },
        )

        return JsonResponse({"success": True, "message": "Marks saved successfully!"})
    return JsonResponse({"error": "Invalid request"}, status=400)

@login_required(login_url='/accounts/login-page/')
def edit_marks_page(request):
    user = request.user
    courses = Course.objects.filter(usercourse__user=user)
    return render(request, "accounts/edit_marks.html", {"username": user.username, "courses": courses})

@login_required(login_url='/accounts/login-page/')
def get_marks_api(request, course_id):
    user = request.user
    try:
        course = Course.objects.get(id=course_id)
        marks = Marks.objects.filter(student=user, course=course).first()
        if marks:
            data = {
                "quiz1": marks.quiz1,
                "quiz2": marks.quiz2,
                "quiz3": marks.quiz3,
                "attendance": marks.attendance,
                "assignment": marks.assignment,
                "presentation": marks.presentation,
                "termexam": marks.termexam,
            }
        else:
            data = {
                "quiz1": 0, "quiz2":0, "quiz3":0,
                "attendance":0, "assignment":0,
                "presentation":0, "termexam":0
            }
        return JsonResponse({"success": True, "marks": data})
    except Course.DoesNotExist:
        return JsonResponse({"success": False, "message": "Course not found"})


@login_required(login_url='/accounts/login-page/')
@csrf_exempt
def edit_marks_api(request):
    if request.method == "POST":
        data = json.loads(request.body)
        student = request.user
        course = Course.objects.get(id=data['course_id'])

        marks, created = Marks.objects.update_or_create(
            student=student,
            course=course,
            defaults={
                "quiz1": data.get("quiz1", 0),
                "quiz2": data.get("quiz2", 0),
                "quiz3": data.get("quiz3", 0),
                "attendance": data.get("attendance", 0),
                "assignment": data.get("assignment", 0),
                "presentation": data.get("presentation", 0),
                "termexam": data.get("termexam", 0),
            },
        )
        return JsonResponse({"success": True, "message": "Marks updated successfully!"})
    return JsonResponse({"success": False, "message": "Invalid request"})

@login_required(login_url='/accounts/login-page/')
def upload_syllabus_page(request):
    user_courses = Course.objects.filter(usercourse__user=request.user)

    if request.method == "POST":
        course_id = request.POST.get("course")
        topic_name = request.POST.get("topic_name")
        lecture_slide = request.FILES.get("lecture_slide")

        if course_id and topic_name and lecture_slide:
            course = Course.objects.get(id=course_id)
            Syllabus.objects.create(
                user=request.user,
                course=course,
                topic_name=topic_name,
                lecture_slide=lecture_slide
            )
            return redirect('upload_syllabus_page')  

    return render(request, "accounts/upload_syllabus.html", {"courses": user_courses})

@login_required(login_url='/accounts/login-page/')
def upload_syllabus(request):
    user_courses = UserCourse.objects.filter(user=request.user).values_list('course', flat=True)
    courses = Course.objects.filter(id__in=user_courses)

    if request.method == 'POST':
        course_id = request.POST.get('course')
        topic_name = request.POST.get('topic_name')
        lecture_slide = request.FILES.get('lecture_slide')

        if not course_id or not topic_name or not lecture_slide:
            messages.error(request, "All fields are required!")
        else:
            course = Course.objects.get(id=course_id)
            Syllabus.objects.create(
                user=request.user,
                course=course,
                topic_name=topic_name,
                lecture_slide=lecture_slide
            )
            messages.success(request, "Syllabus uploaded successfully!")
            return redirect('upload_syllabus_page')

    return render(request, 'accounts/upload_syllabus.html', {'courses': courses})

@login_required(login_url='/accounts/login-page/')
def selected_syllabus(request):
    syllabi = Syllabus.objects.filter(user=request.user).order_by('-uploaded_at')

    for s in syllabi:
        if s.lecture_slide:
            full_name = s.lecture_slide.name         
            filename = full_name.split('/')[-1]   
            ext = filename.split('.')[-1].upper() 

            s.display_name = filename
            s.file_type = ext
        else:
            s.display_name = ""
            s.file_type = ""

    return render(request, 'accounts/selected_syllabus.html', {
        'syllabi': syllabi
    })

@login_required(login_url='/accounts/login-page/')
@require_POST
def delete_syllabus(request, pk):
    s = get_object_or_404(Syllabus, pk=pk, user=request.user)

    if s.lecture_slide:
        s.lecture_slide.delete(save=False)
    s.delete()
    messages.success(request, "Syllabus deleted.")
    return redirect('selected_syllabus_page')

@login_required(login_url='/accounts/login-page/')
def view_syllabus(request, pk):
    syllabus = get_object_or_404(Syllabus, pk=pk, user=request.user)
    file_path = syllabus.lecture_slide.path
    file_mime, _ = mimetypes.guess_type(file_path)
    response = FileResponse(open(file_path, 'rb'), content_type=file_mime)
    response['Content-Disposition'] = f'inline; filename="{syllabus.lecture_slide.name.split("/")[-1]}"'
    return response

@login_required(login_url='/accounts/login-page/')
def calculate_cgpa_page(request):
    return render(request, 'accounts/calculate_cgpa.html')

# overall performance
'''
@login_required(login_url='/accounts/login-page/')
def overall_performance(request):
    user = request.user

    student_marks = Marks.objects.filter(student=user)

    if not student_marks.exists():
        context = {
            "completed_quizzes": 0,
            "performance_rate": 0,
            "improvement_rate": 0,
            "predicted_performance": 0,
            "courses": [],
            "attendance_by_course": {},
        }
        return render(request, "accounts/overall_performance.html", context)

    completed_quizzes = 0
    for mark in student_marks:
        completed_quizzes += sum(1 for q in [mark.quiz1, mark.quiz2, mark.quiz3] if q > 0)

    avg_total_score = (
        student_marks.aggregate(
            avg=Avg(
                (F('quiz1') + F('quiz2') + F('quiz3') +
                 F('attendance') + F('assignment') + F('presentation') + F('termexam')) / 7
            )
        )["avg"] or 0
    )
    performance_rate = round(avg_total_score, 2)

    marks_sorted = student_marks.order_by('uploaded_at')
    if marks_sorted.count() >= 2:
        first_score = (marks_sorted.first().quiz1 + marks_sorted.first().quiz2 + marks_sorted.first().quiz3) / 3
        last_score = (marks_sorted.last().quiz1 + marks_sorted.last().quiz2 + marks_sorted.last().quiz3) / 3
        improvement_rate = round(last_score - first_score, 2)
    else:
        improvement_rate = 0

    courses = []
    attendance_by_course = {}

    for mark in student_marks:
        total = (mark.quiz1 + mark.quiz2 + mark.quiz3 +
                 mark.attendance + mark.assignment + mark.presentation + mark.termexam) / 7
        courses.append({
            "name": mark.course.name,
            "performance": round(total, 2)
        })

        course_attendance = QuizAttendance.objects.filter(student=user, course=mark.course)
        attendance_by_course[mark.course.name] = {}

        for week in range(1, 5):
            attended = course_attendance.filter(week=week, attended=True).exists()
            attendance_by_course[mark.course.name][f"Week {week}"] = "✔️" if attended else "❌"

    if performance_rate < 60:
        predicted_performance = round(performance_rate + 10, 2)
    elif performance_rate < 80:
        predicted_performance = round(performance_rate + 5, 2)
    else:
        predicted_performance = round(performance_rate + 2, 2)

    courses_json = json.dumps(courses)
    context = {
        "completed_quizzes": completed_quizzes,
        "performance_rate": performance_rate,
        "improvement_rate": improvement_rate,
        "predicted_performance": predicted_performance,
        "courses_json": courses_json,  
        "attendance_by_course": attendance_by_course,
    }
    return render(request, "accounts/overall_performance.html", context)
'''
#------------ai overall feedback + performance-----------
def generate_overall_ai_feedback(courses, performance_rate):
    try:
        course_lines = "\n".join([f"{c['name']}: {c['performance']}%" for c in courses])

        prompt = f"""
You are an academic performance advisor.

Student overall performance: {performance_rate}%

Course-wise performance data:
{course_lines}

Tasks:
1. Identify courses that need the most attention (below 60%).
2. Give short, practical improvement advice.
3. Predict achievable performance percentage if the student improves steadily.
4. End with a short motivational line.

Rules:
- Keep it concise.
- No paragraphs.
- Use short bullet points.
- Friendly student tone.
- Do not mention AI.
Format exactly like this:

Focus Areas:
- ...

Improvement Advice:
- ...

Prediction:
- ...

Motivation:
- ...
"""

        model = genai.GenerativeModel("gemini-2.5-flash-lite-preview-09-2025")
        response = model.generate_content(prompt)

        try:
            return response.text.strip()
        except AttributeError:
            return str(response).strip()

    except Exception:
        return (
            "Focus Areas:\n"
            "- Overall data insufficient.\n\n"
            "Improvement Advice:\n"
            "- Keep practicing consistently.\n\n"
            "Prediction:\n"
            "- Gradual improvement expected.\n\n"
            "Motivation:\n"
            "- Keep going strong."
        )

from datetime import timedelta
from django.utils.timezone import now
from django.db.models import Avg, F
from django.contrib.auth.decorators import login_required
import json

@login_required(login_url='/accounts/login-page/')
def overall_performance(request):
    user = request.user

    student_marks = Marks.objects.filter(student=user)

    if not student_marks.exists():
        context = {
            "student": user,
            "completed_quizzes": 0,
            "performance_rate": 0,
            "improvement_rate": 0,
            "predicted_performance": 0,
            "courses": [],
            "courses_json": "[]",
            "attendance_by_course": {},
            "attendance_json": "{}",
            "ai_overall_feedback": "No AI feedback available yet. Add some marks first.",
        }
        return render(request, "accounts/overall_performance.html", context)

    completed_quizzes = 0
    for mark in student_marks:
        completed_quizzes += sum(
            1 for q in [mark.quiz1, mark.quiz2, mark.quiz3] if q > 0
        )

    avg_total_score = (
        student_marks.aggregate(
            avg=Avg(
                (F('quiz1') + F('quiz2') + F('quiz3') +
                 F('attendance') + F('assignment') +
                 F('presentation') + F('termexam')) / 7
            )
        )["avg"] or 0
    )
    performance_rate = round(avg_total_score, 2)

    marks_sorted = student_marks.order_by('uploaded_at')
    if marks_sorted.count() >= 2:
        first_score = (
            marks_sorted.first().quiz1 +
            marks_sorted.first().quiz2 +
            marks_sorted.first().quiz3
        ) / 3
        last_score = (
            marks_sorted.last().quiz1 +
            marks_sorted.last().quiz2 +
            marks_sorted.last().quiz3
        ) / 3
        improvement_rate = round(last_score - first_score, 2)
    else:
        improvement_rate = 0

    courses = []
    attendance_by_course = {}

    # ---- FIXED PART STARTS HERE ----

    all_courses = Course.objects.all()

    for course in all_courses:
        attendance_by_course[course.name] = {
            "Week 1": 0,
            "Week 2": 0,
            "Week 3": 0,
            "Week 4": 0,
        }

        attempts = AIQuizAttempt.objects.filter(
            student=user,
            course=course
        )

        for attempt in attempts:
            day = attempt.created_at.day

            if day <= 7:
                week = "Week 1"
            elif day <= 14:
                week = "Week 2"
            elif day <= 21:
                week = "Week 3"
            else:
                week = "Week 4"

            attendance_by_course[course.name][week] = 1

    # ---- FIXED PART ENDS HERE ----

    for mark in student_marks:
        total = (
            mark.quiz1 + mark.quiz2 + mark.quiz3 +
            mark.attendance + mark.assignment +
            mark.presentation + mark.termexam
        ) / 7

        courses.append({
            "name": mark.course.name,
            "performance": round(total, 2)
        })

    try:
        ai_overall_feedback = generate_overall_ai_feedback(
            courses, performance_rate
        ) or "No AI feedback available."
    except Exception:
        ai_overall_feedback = "No AI feedback available."

    if performance_rate < 60:
        predicted_performance = round(performance_rate + 10, 2)
    elif performance_rate < 80:
        predicted_performance = round(performance_rate + 5, 2)
    else:
        predicted_performance = round(performance_rate + 2, 2)

    courses_json = json.dumps(courses)
    attendance_json = json.dumps(attendance_by_course)

    context = {
        "student": user,
        "completed_quizzes": completed_quizzes,
        "performance_rate": performance_rate,
        "improvement_rate": improvement_rate,
        "predicted_performance": predicted_performance,
        "courses": courses,
        "courses_json": courses_json,
        "attendance_by_course": attendance_by_course,
        "attendance_json": attendance_json,
        "ai_overall_feedback": ai_overall_feedback,
    }

    return render(request, "accounts/overall_performance.html", context)

# AIUB grade points mapping
GRADE_POINTS = {
    "A+": 4.00,
    "A": 3.75,
    "B+": 3.50,
    "B": 3.25,
    "C+": 3.00,
    "C": 2.75,
    "D+": 2.50,
    "D": 2.25,
    "F": 0.00
}

@login_required(login_url='/accounts/login-page/')
def current_semester_cg(request):
    courses = []
    cg = None
    total_credits = None

    if request.method == "POST":
        num_courses = int(request.POST.get("num_courses", 0))
        total_points = 0
        total_credits = 0

        for i in range(1, num_courses + 1):
            grade = request.POST.get(f"grade_{i}")
            credits = float(request.POST.get(f"credits_{i}", 0))
            if grade in GRADE_POINTS and credits > 0:
                total_points += GRADE_POINTS[grade] * credits
                total_credits += credits

        if total_credits > 0:
            cg = round(total_points / total_credits, 2)

    return render(request, "accounts/current_semester_cg.html", {
        "courses": courses,
        "cg": cg,
        "total_credits": total_credits,
        "grades": list(GRADE_POINTS.keys())
    })

# -------------- File extraction functions --------------------
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

def extract_text_from_pdf(file_path):
    """
    Extract text from a PDF file.
    Handles both plain text PDFs and scanned image PDFs using OCR.
    """
    text = ""

    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                else:
                    pil_image = page.to_image(resolution=300).original
                    ocr_text = pytesseract.image_to_string(pil_image)
                    text += ocr_text + "\n"
    except Exception as e:
        text += f"\n[Error reading PDF {os.path.basename(file_path)}: {str(e)}]\n"

    return text

def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
    except Exception as e:
        text += f"\n[Error reading DOCX {os.path.basename(file_path)}: {str(e)}]\n"
    return text

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text += cell.text + "\n"
    except Exception as e:
        text += f"\n[Error reading PPTX {os.path.basename(file_path)}: {str(e)}]\n"
    return text

def extract_text_from_excel(file_path):
    text = ""
    wb = openpyxl.load_workbook(file_path)
    
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                text += "\t".join(row_text) + "\n"
 
        for image in sheet._images:
            img_bytes = image._data() 
            pil_image = Image.open(BytesIO(img_bytes))
            ocr_text = pytesseract.image_to_string(pil_image)
            if ocr_text.strip():
                text += "\n[Image Text]: " + ocr_text + "\n"

    return text

def extract_text_from_image(file_path):
    """
    Extract text from an image or screenshot using OCR.
    Supports PNG, JPG, JPEG, BMP, etc.
    """
    text = ""
    try:
        img = Image.open(file_path)
        ocr_text = pytesseract.image_to_string(img)
        if ocr_text.strip():
            text += ocr_text
        else:
            text += "[No text detected in image]"
    except Exception as e:
        text += f"\n[Error reading image {os.path.basename(file_path)}: {str(e)}]\n"
    return text
# -------------------- Aggregate all syllabus content --------------------
def get_syllabus_text(user):
    """Return all text from all syllabus files for this user."""
    syllabi = Syllabus.objects.filter(user=user)
    full_text = ""
    for s in syllabi:
        path = s.lecture_slide.path
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            full_text += extract_text_from_pdf(path) + "\n"
        elif ext == ".docx":
            full_text += extract_text_from_docx(path) + "\n"
        elif ext == ".pptx":
            full_text += extract_text_from_pptx(path) + "\n"
        elif ext == ".xlsx": 
            full_text += extract_text_from_excel(path) + "\n"
        elif ext in [".png", ".jpg", ".jpeg", ".bmp"]:
            full_text += extract_text_from_image(path) + "\n"
        else:
            full_text += f"\n[Unsupported file type {s.lecture_slide.name} skipped]\n"
    return full_text

# -------------------- AI Assistant view --------------------
@login_required(login_url='/accounts/login-page/')
def ai_assistant(request):
    user_prompt = ""

    if 'conversation' not in request.session:
        request.session['conversation'] = []

    if request.method == "POST":
        user_prompt = request.POST.get("prompt", "").strip()
        if user_prompt:
            conversation_history = request.session.get('conversation', [])

            syllabus_text = get_syllabus_text(request.user)
            full_prompt = f"""You are an AI tutor with full general knowledge and syllabus awareness.  
Use the syllabus text below when the user's question is related to the course.  
If the question is NOT related to the syllabus, answer normally using your full capabilities.
Syllabus content:
{syllabus_text}
Conversation so far:
"""
            for msg in conversation_history:
                full_prompt += f"{msg['role']}: {msg['content']}\n"
            full_prompt += f"user: {user_prompt}"

            try:
                model = genai.GenerativeModel("gemini-2.5-flash-lite-preview-09-2025")
                result = model.generate_content(full_prompt)
                ai_raw_response = result.text
                ai_response = markdown2.markdown(
                    ai_raw_response,
                    extras=["fenced-code-blocks", "break-on-newline"]
                )

            except Exception as e:
                ai_response = f"⚠️ Error generating AI response: {str(e)}"

            conversation_history.append({'role': 'user', 'content': user_prompt})
            conversation_history.append({'role': 'ai', 'content': ai_response})
            request.session['conversation'] = conversation_history

    return render(request, "accounts/ai_assistant.html", {
        "conversation": request.session.get('conversation', []),
        "prompt": user_prompt
    })

#-------------------
def tools_page(request):
    user = request.user
    return render(request, "accounts/tools_page.html",{
        "username": user.username
    })
#-------------------------------
def get_course_syllabus_text(course, user, filename=None):
    syllabi = Syllabus.objects.filter(course=course, user=user)
    if filename:
        syllabi = syllabi.filter(lecture_slide__icontains=filename)

    text = ""
    for s in syllabi:
        if not s.lecture_slide:
            continue
        path = s.lecture_slide.path
        if not os.path.exists(path):
            continue
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            text += extract_text_from_pdf(path)
        elif ext == ".docx":
            text += extract_text_from_docx(path)
        elif ext == ".pptx":
            text += extract_text_from_pptx(path)
        elif ext in [".png", ".jpg", ".jpeg"]:
            text += extract_text_from_image(path)
        text += "\n"

    return text[:12000]

# AI Quiz Page
@login_required(login_url='/accounts/login-page/')
def ai_quiz_page(request):
    courses = Course.objects.filter(usercourse__user=request.user)
    topics = Syllabus.objects.filter(user=request.user).values_list("lecture_slide", flat=True).distinct()
    topics = [os.path.basename(t) for t in topics if t]

    return render(request, "accounts/ai_quiz.html", {
        "courses": courses,
        "topics": topics
    })

# Generate AI Quiz
@login_required(login_url='/accounts/login-page/')
def generate_ai_quiz(request):
    if request.method != "POST":
        return JsonResponse({"status": "error", "message": "Invalid request method"})

    data = json.loads(request.body)
    course_id = data.get("course_id")
    filename = data.get("topic_name")
    instruction = data.get("instruction", "").strip()

    if not course_id:
        return JsonResponse({"status": "error", "message": "Course not selected"})

    try:
        course = Course.objects.get(
            id=int(course_id),
            usercourse__user=request.user
        )
    except Course.DoesNotExist:
        return JsonResponse({"status": "error", "message": "Course not found"})

    syllabus_text = get_course_syllabus_text(course, request.user, filename)

    if not syllabus_text.strip():
        return JsonResponse({"status": "error", "message": "No syllabus content found for this topic"})

    # instruction is REQUIRED for controlling output
    if not instruction:
        return JsonResponse({
            "status": "error",
            "message": "Please specify what you want the AI to generate"
        })

    prompt = f"""
Follow the user's instruction EXACTLY.
Especially the number of questions requested.

Each MCQ must include:
- question_text
- options a, b, c, d
- correct_answer (a/b/c/d)
- marks = 1

Return ONLY a valid JSON array.
No markdown. No explanation.

USER INSTRUCTION:
{instruction}

SYLLABUS CONTENT:
{syllabus_text}
"""

    try:
        model = genai.GenerativeModel("gemini-2.5-flash-lite-preview-09-2025")
        response = model.generate_content(prompt)

        raw = response.text.strip()
        raw = raw.replace("```json", "").replace("```", "")
        questions = json.loads(raw)

        return JsonResponse({
            "status": "success",
            "questions": questions
        })

    except Exception as e:
        return JsonResponse({
            "status": "error",
            "message": str(e)
        })

    
'''
@login_required(login_url='/accounts/login-page/')
def submit_ai_quiz(request):
    data = json.loads(request.body)
    total = 0
    max_total = 0

    for q in data.get("questions", []):
        marks = q.get("marks", 1)
        max_total += marks
        if q.get("student_answer") == q.get("correct_answer"):
            total += marks

    AIQuizAttempt.objects.create(
        student=request.user,
        course_id=data.get("course_id"),
        topic_name=data.get("topic_name"),
        score=total,
        max_score=max_total
    )

    return JsonResponse({
        "status": "success",
        "total_marks": total,
        "max_marks": max_total
    })
'''
@login_required(login_url='/accounts/login-page/')
def submit_ai_quiz(request):
    data = json.loads(request.body)
    total = 0
    max_total = 0
    detailed_feedback = []
    syllabus_cache = {} 

    for q in data.get("questions", []):
        marks = q.get("marks", 1)
        max_total += marks
        student_answer = q.get("student_answer")
        correct_answer = q.get("correct_answer")
        filename = data.get("topic_name")

        is_correct = student_answer == correct_answer
        if is_correct:
            total += marks

        if filename not in syllabus_cache:
            syllabus_cache[filename] = get_course_syllabus_text(
                Course.objects.get(id=data.get("course_id")),
                user=request.user,
                filename=filename
            )
        syllabus_text = syllabus_cache[filename]

        if not is_correct and student_answer:
            #explanation = f"Review the following key points from the syllabus:\n{syllabus_text[:1000]}."
            explanation = f"You wronged question."
        elif is_correct:
            explanation = f"Answered correctly."
        else:
            explanation = f"You did not answer this question."

        detailed_feedback.append({
            "question_text": q.get("question_text"),
            "student_answer": student_answer,
            "correct_answer": correct_answer,
            "is_correct": is_correct,
            "explanation": explanation
        })

    AIQuizAttempt.objects.create(
        student=request.user,
        course_id=data.get("course_id"),
        topic_name=data.get("topic_name"),
        score=total,
        max_score=max_total
    )

    return JsonResponse({
        "status": "success",
        "total_marks": total,
        "max_marks": max_total,
        "detailed_feedback": detailed_feedback
    })

#progress
@login_required(login_url='/accounts/login-page/')
def student_progress(request):
    student = request.user

    attempts = AIQuizAttempt.objects.filter(
        student=student
    ).order_by('created_at')

    progress_data = defaultdict(lambda: defaultdict(list))

    for attempt in attempts:
        course = attempt.course.name
        topic = attempt.topic_name if attempt.topic_name else "General"

        progress_data[course][topic].append({
            "score": attempt.percentage(),
            "date": attempt.created_at.strftime("%Y-%m-%d")
        })

    context = {
        "username": student.username,
        "progress_json": json.dumps(progress_data)
    }

    return render(request, "accounts/student_progress.html", context)
#-----------------------------TEACHER--------------------------------------------
#-----------------------------TEACHER--------------------------------------------
# teacher Dashboard page (protected)
@login_required(login_url='/accounts/login-page/')
def teacher_dashboard(request):
    user = request.user
    context = {
        "username": user.username,
        "role": getattr(user, "role", "N/A"),
    }
    return render(request, "accounts/teacher_dashboard.html", context)

# ---------- Teacher Course Selection Page ----------
@login_required(login_url='/accounts/login-page/')
def teacher_course_selection_page(request):
    user = request.user
    courses = Course.objects.all()

    selections = TeacherCourseSelection.objects.filter(teacher=user)
    # Dict: { "Spring": ["CSE101", "CSE102"], "Fall": ["CSE201"] }
    semester_courses = {
        sel.semester: list(sel.courses.values_list('code', flat=True))
        for sel in selections
    }
    context = {
        "courses": courses,
        "semester_courses_json": json.dumps(semester_courses),  # for JS
        "username": user.username,
    }
    return render(request, "accounts/teacher_course_selection.html", context)

@csrf_exempt
@login_required(login_url='/accounts/login-page/')
def save_teacher_courses_api(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            semester = data.get("semester")
            selected_codes = data.get("courses", [])

            if not semester or not isinstance(selected_codes, list):
                return JsonResponse({"success": False, "message": "Invalid data"})

            user = request.user
            selection, _ = TeacherCourseSelection.objects.get_or_create(
                teacher=user, semester=semester
            )
            courses = Course.objects.filter(code__in=selected_codes)
            selection.courses.set(courses)
            selection.save()

            return JsonResponse({
                "success": True,
                "message": f"Saved {len(selected_codes)} course(s) for {semester}.",
                "semester": semester,
                "courses": selected_codes
            })

        except Exception as e:
            return JsonResponse({"success": False, "message": str(e)})

    return JsonResponse({"success": False, "message": "Invalid request"})

@login_required(login_url='/accounts/login-page/')
def teacher_selected_courses_page(request):
    user = request.user
    selections = TeacherCourseSelection.objects.filter(teacher=user)
    return render(request, "accounts/teacher_selected_courses.html", {
        "username": user.username,
        "selections": selections
    })
#----------------------
@login_required(login_url='/accounts/login-page/')
def upload_student_info_page(request):
    user = request.user
    selections = TeacherCourseSelection.objects.filter(teacher=user)

    semester_courses_dict = {}
    for sel in selections:
        course_list = list(sel.courses.values_list('name', flat=True))
        semester_courses_dict[sel.semester] = course_list

    if request.method == "POST":
        semester = request.POST.get('semester')
        course_name = request.POST.get('course')
        section_name = request.POST.get('section')
        csv_file = request.FILES.get('studentFile')

        if not (semester and course_name and section_name and csv_file):
            messages.error(request, "All fields are required!")
            return redirect('upload_student_info_page')

        try:
            course_obj = Course.objects.get(name=course_name)
        except Course.DoesNotExist:
            messages.error(request, f"Course '{course_name}' not found!")
            return redirect('upload_student_info_page')

        section_obj, _ = Section.objects.get_or_create(
            teacher=user,
            semester=semester,
            course=course_obj,
            section_name=section_name
        )

        section_obj.students.all().delete()
        try:
            decoded_file = csv_file.read().decode('utf-8')
            io_string = io.StringIO(decoded_file)
            reader = csv.DictReader(io_string)

            count = 0
            for row in reader:
                student_name = row.get('student_name') or row.get('name')
                student_id = row.get('student_id')
                email = row.get('email')

                if not student_name or not student_id or not email:
                    continue  

                StudentInfo.objects.create(
                    section=section_obj,
                    student_name=student_name,
                    student_id=student_id,
                    email=email
                )
                count += 1

            messages.success(
                request,
                f"Successfully replaced and uploaded {count} students for {course_name} ({section_name}) - {semester}."
            )
            return redirect('upload_student_info_page')

        except Exception as e:
            messages.error(request, f"Error processing CSV file: {str(e)}")
            return redirect('upload_student_info_page')

    context = {
        "username": user.username,
        "semester_courses_dict": semester_courses_dict,
        "semester_courses_json": json.dumps(semester_courses_dict),  
        "section_choices": ["A", "B", "C", "D", "E", "F"],
    }
    return render(request, "accounts/upload_student_info.html", context)

@login_required(login_url='/accounts/login-page/')
def uploaded_students_page(request):
    user = request.user

    sections = Section.objects.filter(teacher=user).order_by('semester', 'course__name', 'section_name')

    data = {}
    for section in sections:
        sem = section.semester
        course = section.course.name
        sec_name = section.section_name
        students = section.students.all()

        if sem not in data:
            data[sem] = {}
        if course not in data[sem]:
            data[sem][course] = {}
        data[sem][course][sec_name] = students

    context = {
        "username": user.username,
        "uploaded_data": data
    }

    return render(request, "accounts/uploaded_students.html", context)

@login_required(login_url='/accounts/login-page/')
def delete_section(request):
    if request.method == "POST":
        user = request.user
        semester = request.POST.get('semester')
        course_name = request.POST.get('course')
        section_name = request.POST.get('section')

        try:
            course_obj = Course.objects.get(name=course_name)
            section = Section.objects.get(
                teacher=user,
                semester=semester,
                course=course_obj,
                section_name=section_name
            )
            section.delete()  
            messages.success(request, f"Deleted section {section_name} ({course_name}) successfully!")
        except Section.DoesNotExist:
            messages.error(request, "Section not found or you don't have permission!")

    return redirect('uploaded_students_page')

def get_teacher_semester_courses(user):
    selections = TeacherCourseSelection.objects.filter(teacher=user)
    semester_courses = {}
    for sel in selections:
        course_list = list(sel.courses.values_list('name', flat=True))
        semester_courses[sel.semester] = course_list
    return semester_courses

@login_required(login_url='/accounts/login-page/')
def teacher_post_marks_page(request):
    """
    Teacher page to upload marks via CSV for a selected semester/course/section.
    Automatically replaces old marks for same semester/course/section uploaded by same teacher.
    CSV expected columns:
      student_email or student_name, quiz1, quiz2, quiz3, attendance, assignment, presentation, termexam
    Missing numeric fields default to 0.
    """
    user = request.user
    semester_courses_dict = get_teacher_semester_courses(user)
    section_choices = ["A", "B", "C", "D", "E", "F"]

    if request.method == "POST":
        semester = request.POST.get('semester')
        course_name = request.POST.get('course')
        section_name = request.POST.get('section')
        csv_file = request.FILES.get('marksFile')

        if not (semester and course_name and section_name and csv_file):
            messages.error(request, "All fields are required (semester, course, section, and CSV file).")
            return redirect('teacher_post_marks_page')
        try:
            course_obj = Course.objects.get(name=course_name)
        except Course.DoesNotExist:
            messages.error(request, f"Course '{course_name}' not found.")
            return redirect('teacher_post_marks_page')

        section_obj, _ = Section.objects.get_or_create(
            teacher=user,
            semester=semester,
            course=course_obj,
            section_name=section_name
        )

        try:
            decoded = csv_file.read().decode('utf-8')
            io_string = io.StringIO(decoded)
            reader = csv.DictReader(io_string)

            inserted = 0
            updated = 0
            from django.contrib.auth import get_user_model
            User = get_user_model()
            Marks.objects.filter(
                course=course_obj,
                semester=semester,
                section=section_name,
                course__section__teacher=user
            ).delete()

            with transaction.atomic():
                for row in reader:
                    email = (row.get('student_email') or row.get('email') or "").strip().lower()
                    student_name = (row.get('student_name') or row.get('name') or "").strip()

                    student_user = None
                    student_info = None

                    if email:
                        try:
                            student_user = User.objects.get(email=email)
                            student_info = StudentInfo.objects.filter(email=email, section=section_obj).first()
                        except User.DoesNotExist:
                            student_user = None

                    if not student_user and student_name:
                        student_info = StudentInfo.objects.filter(
                            student_name__iexact=student_name,
                            section=section_obj
                        ).first()
                        if student_info and student_info.email:
                            student_user = User.objects.filter(email=student_info.email).first()

                    if not student_user:
                        continue 
                    def to_float(v):
                        try:
                            return float(v)
                        except Exception:
                            return 0.0

                    quiz1 = to_float(row.get('quiz1') or row.get('quiz_1') or row.get('q1'))
                    quiz2 = to_float(row.get('quiz2') or row.get('quiz_2') or row.get('q2'))
                    quiz3 = to_float(row.get('quiz3') or row.get('quiz_3') or row.get('q3'))
                    attendance = to_float(row.get('attendance') or row.get('att'))
                    assignment = to_float(row.get('assignment') or row.get('assign'))
                    presentation = to_float(row.get('presentation') or row.get('presentation_score') or row.get('present'))
                    termexam = to_float(row.get('termexam') or row.get('term_exam') or row.get('final'))

                    Marks.objects.create(
                        student=student_user,
                        course=course_obj,
                        semester=semester,
                        section=section_name,
                        quiz1=quiz1,
                        quiz2=quiz2,
                        quiz3=quiz3,
                        attendance=attendance,
                        assignment=assignment,
                        presentation=presentation,
                        termexam=termexam,
                        published=False
                    )
                    inserted += 1

            messages.success(
                request,
                f"Marks uploaded successfully for {course_name} (Sem {semester}, Sec {section_name}). "
                f"Old marks replaced with {inserted} new entries."
            )
            return redirect('teacher_post_marks_page')

        except Exception as e:
            messages.error(request, f"Error processing CSV: {str(e)}")
            return redirect('teacher_post_marks_page')

    context = {
        "username": user.username,
        "semester_courses_dict": semester_courses_dict,
        "semester_courses_json": json.dumps(semester_courses_dict),
        "section_choices": section_choices,
    }
    return render(request, 'accounts/teacher_post_marks.html', context)

@login_required(login_url='/accounts/login-page/')
def publish_marks(request):
    """
    Publish marks for a given semester/course/section (POST).
    This will set published=True for all Marks objects for that course and for students in that section.
    """
    if request.method != "POST":
        return redirect('teacher_post_marks_page')

    user = request.user
    semester = request.POST.get('semester')
    course_name = request.POST.get('course')
    section_name = request.POST.get('section')

    if not (semester and course_name and section_name):
        messages.error(request, "Semester, course and section required to publish marks.")
        return redirect('teacher_post_marks_page')

    try:
        course_obj = Course.objects.get(name=course_name)
    except Course.DoesNotExist:
        messages.error(request, "Course not found.")
        return redirect('teacher_post_marks_page')

    try:
        section_obj = Section.objects.get(teacher=user, semester=semester, course=course_obj, section_name=section_name)
    except Section.DoesNotExist:
        messages.error(request, "Section not found.")
        return redirect('teacher_post_marks_page')

    student_infos = StudentInfo.objects.filter(section=section_obj)
    user_emails = [s.email for s in student_infos if s.email]

    from django.contrib.auth import get_user_model
    User = get_user_model()
    students_qs = User.objects.filter(email__in=user_emails)
    marks_qs = Marks.objects.filter(course=course_obj, student__in=students_qs)
    updated = marks_qs.update(published=True)

    messages.success(request, f"Published marks for {section_name} ({course_name}) — {updated} mark rows published.")
    return redirect('teacher_post_marks_page')

@login_required(login_url='/accounts/login-page/')
def student_my_marks_page(request):
    """
    Student view: show only published Marks for the logged-in student.
    """
    user = request.user
    if hasattr(user, 'role') and user.role != 'Student':
        pass

    marks = Marks.objects.filter(student=user, published=True).select_related('course').order_by('course__name')

    return render(request, 'accounts/student_my_marks.html', {
        'username': user.username,
        'marks': marks
    })

@login_required(login_url='/accounts/login-page/')
def teacher_uploaded_marks_page(request):
    """
    Show marks uploaded/published by the teacher, filtered by semester, course, and section.
    Uses semester and section stored in Marks model itself.
    """
    user = request.user
    teacher_marks = Marks.objects.filter(course__section__teacher=user)

    semesters = teacher_marks.values_list('semester', flat=True).distinct()
    courses = Course.objects.filter(id__in=teacher_marks.values_list('course', flat=True).distinct())
    sections = teacher_marks.values_list('section', flat=True).distinct()

    selected_sem = request.GET.get('semester')
    selected_course_id = request.GET.get('course')
    selected_section = request.GET.get('section')

    marks = Marks.objects.none()  
    selected_course_name = None

    if selected_course_id:
        try:
            course_obj = Course.objects.get(id=selected_course_id)
            selected_course_name = course_obj.name

            if selected_sem and selected_section:
                marks = teacher_marks.filter(
                    course=course_obj,
                    semester=selected_sem,
                    section=selected_section
                ).select_related('student', 'course')

        except Course.DoesNotExist:
            selected_course_name = None
            marks = Marks.objects.none()

    context = {
        'username': user.username,
        'semesters': semesters,
        'courses': courses,
        'sections': sections,
        'marks': marks,
        'selected_sem': selected_sem,
        'selected_course_id': int(selected_course_id) if selected_course_id else None,
        'selected_section': selected_section,
        'selected_course_name': selected_course_name,
    }

    return render(request, 'accounts/teacher_published_marks.html', context)
#
@login_required(login_url='/accounts/login-page/')
def teacher_view_performance(request):
    user = request.user
    sections = Section.objects.filter(teacher=user).order_by(
        'semester', 'course__name', 'section_name'
    )
    students_map = {}

    for section in sections:
        for student_info in section.students.all():

            student_user = User.objects.filter(email=student_info.email).first()
            if not student_user:
                continue

            student_marks = Marks.objects.filter(student=student_user).select_related(
                'course'
            ).order_by('uploaded_at')

            completed_quizzes = 0
            performance_rate = 0
            improvement_rate = 0
            predicted_performance = 0
            courses_list = []
            attendance_by_course = {}

            if student_marks.exists():
                completed_quizzes = sum(
                    1
                    for m in student_marks
                    for q in [m.quiz1, m.quiz2, m.quiz3]
                    if (q is not None and q > 0)
                )
                totals = []
                for m in student_marks:
                    row_total = (
                        (m.quiz1 or 0)
                        + (m.quiz2 or 0)
                        + (m.quiz3 or 0)
                        + (m.attendance or 0)
                        + (m.assignment or 0)
                        + (m.presentation or 0)
                        + (m.termexam or 0)
                    ) / 7
                    totals.append(row_total)

                performance_rate = round(sum(totals) / len(totals), 2) if totals else 0
                if student_marks.count() >= 2:
                    first = student_marks.first()
                    last = student_marks.last()
                    first_avg = (
                        (first.quiz1 or 0)
                        + (first.quiz2 or 0)
                        + (first.quiz3 or 0)
                    ) / 3
                    last_avg = (
                        (last.quiz1 or 0)
                        + (last.quiz2 or 0)
                        + (last.quiz3 or 0)
                    ) / 3
                    improvement_rate = round(last_avg - first_avg, 2)

                if performance_rate < 60:
                    predicted_performance = round(performance_rate + 10, 2)
                elif performance_rate < 80:
                    predicted_performance = round(performance_rate + 5, 2)
                else:
                    predicted_performance = round(performance_rate + 2, 2)

                course_map = {}
                for m in student_marks:
                    cname = m.course.name if m.course else "Unknown"
                    row_total = (
                        (m.quiz1 or 0)
                        + (m.quiz2 or 0)
                        + (m.quiz3 or 0)
                        + (m.attendance or 0)
                        + (m.assignment or 0)
                        + (m.presentation or 0)
                        + (m.termexam or 0)
                    ) / 7
                    course_map.setdefault(cname, []).append(row_total)

                for cname, vals in course_map.items():
                    avg_course = round(sum(vals) / len(vals), 2)
                    courses_list.append({"name": cname, "performance": avg_course})

                for course_info in courses_list:
                    ca = QuizAttendance.objects.filter(
                        student=student_user, course__name=course_info["name"]
                    )
                    attendance_weeks = {}
                    for week in range(1, 5):
                        attended = ca.filter(week=week, attended=True).exists()
                        attendance_weeks[f"Week {week}"] = (
                            "✔️" if attended else "❌"
                        )
                    attendance_by_course[course_info["name"]] = attendance_weeks

            students_map[student_info.email] = {
                "name": student_info.student_name
                or student_user.get_full_name()
                or student_user.username,
                "email": student_info.email,
                "section": section.section_name,
                "semester": section.semester,
                "course": section.course.name,
                "total": performance_rate,
                "predicted": predicted_performance,
                "completed_quizzes": completed_quizzes,
                "improvement_rate": improvement_rate,
                "courses": courses_list,
                "attendance_by_course": attendance_by_course,
            }

    students_data = list(students_map.values())
    students_data = sorted(students_data, key=lambda x: x["name"].lower())

    context = {
        "username": user.username,
        "students_data": students_data,
    }

    return render(request, "accounts/teacher_view_performance.html", context)

#----------------------
@login_required(login_url='/accounts/login-page/')
def teacher_upload_syllabus(request):
    teacher = request.user

    semester_courses_qs = TeacherCourseSelection.objects.filter(teacher=teacher).prefetch_related('courses')
    semesters = semester_courses_qs.values_list('semester', flat=True).distinct()

    courses = set()
    for sc in semester_courses_qs:
        for c in sc.courses.all():
            courses.add((c.id, c.name))  
    courses = list(courses)

    sections = Section.objects.filter(teacher=teacher).values_list('section_name', flat=True).distinct()

    if request.method == "POST":
        semester = request.POST.get("semester")
        course_id = request.POST.get("course")
        section = request.POST.get("section")
        syllabus_file = request.FILES.get("syllabus_file")

        if not (semester and course_id and section and syllabus_file):
            messages.error(request, "All fields are required!")
            return redirect("teacher_upload_syllabus")

        existing = TeacherSyllabusUpload.objects.filter(
            teacher=teacher,
            semester=semester,
            course_id=course_id,
            section=section
        ).first()

        if existing:
            existing.syllabus_file.delete()
            existing.delete()

        TeacherSyllabusUpload.objects.create(
            teacher=teacher,
            semester=semester,
            course_id=course_id,
            section=section,
            syllabus_file=syllabus_file
        )

        messages.success(request, "Syllabus uploaded successfully!")
        return redirect("teacher_upload_syllabus")

    context = {
        "username": teacher.username,
        "semesters": semesters,
        "courses": courses,
        "sections": sections,
    }
    return render(request, "accounts/teacher_upload_syllabus.html", context)

#
@login_required(login_url='/accounts/login-page/')
def student_view_syllabus(request):
    user = request.user
    student_courses = UserCourse.objects.filter(user=user).values_list("course", flat=True)

    syllabus_uploads = TeacherSyllabusUpload.objects.filter(course__id__in=student_courses)

    return render(request, "accounts/student_view_syllabus.html", {
        "syllabus": syllabus_uploads,
        "username": user.username
    })
